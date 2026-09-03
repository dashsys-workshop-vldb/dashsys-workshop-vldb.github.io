#!/usr/bin/env python3
"""
Create PowerPoint presentation for DASHSys 2026 Welcome Address
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path

# Initialize presentation (16:9 widescreen)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
COLORS = {
    'dark_blue': RGBColor(26, 35, 50),      # Title text
    'blue': RGBColor(37, 99, 235),          # Data Mgmt
    'purple': RGBColor(124, 58, 237),       # Agentic
    'red': RGBColor(220, 38, 38),           # Evaluation
    'green': RGBColor(5, 150, 105),         # HITL
    'orange': RGBColor(245, 158, 11),       # Infrastructure
    'white': RGBColor(255, 255, 255),
    'gray': RGBColor(107, 114, 128),
}

def add_title_slide(title, subtitle=""):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark_blue']

    # Title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(11.333)
    height = Inches(1.5)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['white']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        left = Inches(1)
        top = Inches(4.2)
        width = Inches(11.333)
        height = Inches(1)
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(28)
        subtitle_frame.paragraphs[0].font.color.rgb = COLORS['white']
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(title, content_lines):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # White background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']

    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(12.333)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['dark_blue']

    # Content
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(11.333)
    height = Inches(5.5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, line in enumerate(content_lines):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = line
        p.font.size = Pt(24)
        p.space_before = Pt(12)
        p.level = 0

    return slide

def add_image_slide(title, image_path, caption=""):
    """Add slide with image"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # White background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']

    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(12.333)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['dark_blue']

    # Image (if exists)
    if Path(image_path).exists():
        left = Inches(1.5)
        top = Inches(1.5)
        slide.shapes.add_picture(str(image_path), left, top, height=Inches(5))

    # Caption
    if caption:
        left = Inches(1)
        top = Inches(6.8)
        width = Inches(11.333)
        height = Inches(0.5)
        caption_box = slide.shapes.add_textbox(left, top, width, height)
        caption_frame = caption_box.text_frame
        caption_frame.text = caption
        caption_frame.paragraphs[0].font.size = Pt(18)
        caption_frame.paragraphs[0].font.color.rgb = COLORS['gray']
        caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

# ============================================================================
# SLIDE 1: Welcome
# ============================================================================
print("Creating Slide 1: Welcome...")
slide1 = add_title_slide(
    "DASHSys 2026",
    "Workshop on Data-AI Systems\nVLDB 2026 • Boston, MA"
)

# ============================================================================
# SLIDE 2: The Convergence Story
# ============================================================================
print("Creating Slide 2: Convergence Story...")
slide2 = add_content_slide(
    "The Convergence Story: How DASHSys Was Born",
    [
        "🔵 DAIS - Data and AI Systems",
        "🟣 DAAS - Data Management for Agentic AI Systems",
        "🟢 DASH - Data Systems with Human-in-the-loop AI",
        "",
        "➡️  Three workshops with overlapping interests",
        "➡️  VLDB asked: Why not merge?",
        "➡️  Result: DASHSys - One unified community",
        "",
        "💡 The convergence itself is the message:",
        "   This is where the field is heading."
    ]
)

# ============================================================================
# SLIDE 3: By the Numbers
# ============================================================================
print("Creating Slide 3: Statistics...")
slide3 = add_content_slide(
    "By the Numbers: Strong Response",
    [
        "📊 SUBMISSIONS",
        "   • 35 total papers submitted",
        "   • Regular Track: 23 papers",
        "   • System Track: 12 papers",
        "",
        "✅ ACCEPTED",
        "   • 12 regular papers (40% acceptance rate)",
        "     - 4 long papers (8 minutes each)",
        "     - 8 short papers (4 minutes each)",
        "   • 2 system track winners",
        "",
        "🏆 COMPETITION WINNERS",
        "   • 1st Place: IBM Research Zurich",
        "   • 2nd Place: Princeton University"
    ]
)

# ============================================================================
# SLIDE 4: Today's Program
# ============================================================================
print("Creating Slide 4: Program...")
slide4 = add_content_slide(
    "Today's Program",
    [
        "09:00-09:10  Welcome & Opening (now!)",
        "09:10-10:30  Keynote Session 1 (Eric Zhu, Fatma Özcan)",
        "10:30-11:00  ☕ Coffee Break + Posters",
        "11:00-12:30  📄 Long Papers (4 × 8 min)",
        "12:30-14:00  🍽️  Lunch Break",
        "14:00-15:30  Keynote Session 2 (Omar Khattab, Eugene Wu)",
        "15:30-16:00  ☕ Coffee Break + Posters",
        "16:00-17:15  📄 Short Papers (8 × 4 min)",
        "17:15-18:00  💬 Panel Discussion + Closing",
        "",
        "🎯 All 12 papers + 2 systems demos at poster sessions!"
    ]
)

# ============================================================================
# SLIDE 5: Paper Themes
# ============================================================================
print("Creating Slide 5: Paper Themes...")
slide5 = add_content_slide(
    "Four Core Research Themes",
    [
        "🔵 42% Data Management for Agentic Systems",
        "   Building the infrastructure agents need",
        "   (GitLake, Data Canvas, Structured State Management)",
        "",
        "🟣 25% Agentic Systems for Data Management",
        "   Agents that help manage and discover data",
        "   (ASMR, GUIDE, Multi-agent Framework)",
        "",
        "🟢 42% Human-Centered & Human-in-the-Loop",
        "   Keeping humans in control",
        "   (Walk Before You Run, AgentTrails, Be Fair!)",
        "",
        "🔴 8% Evaluation, Reliability & Learning",
        "   How to measure and trust these systems (SANA)",
        "   ⚠️  Biggest gap = biggest opportunity!"
    ]
)

# ============================================================================
# SLIDE 6: Cross-Pollination (with image)
# ============================================================================
print("Creating Slide 6: Cross-Pollination...")
img_path = "visualizations/theme_cooccurrence_heatmap.png"
slide6 = add_image_slide(
    "The Hidden Structure: How Themes Connect",
    img_path,
    "Theme-theme co-occurrence: Darker = more papers bridge both themes"
)

# Add key insights as text
left = Inches(8)
top = Inches(1.8)
width = Inches(5)
height = Inches(5)
text_box = slide6.shapes.add_textbox(left, top, width, height)
tf = text_box.text_frame
tf.word_wrap = True

insights = [
    "Strongest Bridges:",
    "• Data ↔ Infrastructure (8×)",
    "• HITL ↔ Provenance (6×)",
    "• Agentic ↔ Multi-Agent (6×)",
    "",
    "Universal Connectors:",
    "⭐ HITL: 11 connections",
    "⭐ Provenance: 10 connections",
    "",
    "⚠️  Critical Gap:",
    "Data ↔ Agentic (0×)",
    "The missing bridge!"
]

for i, line in enumerate(insights):
    if i > 0:
        tf.add_paragraph()
    p = tf.paragraphs[i]
    p.text = line
    p.font.size = Pt(20)
    if line.startswith("⭐") or line.startswith("⚠️"):
        p.font.bold = True

# ============================================================================
# SLIDE 7: Keynotes
# ============================================================================
print("Creating Slide 7: Keynotes...")
slide7 = add_content_slide(
    "Seven Keynote Speakers: Industry Meets Academia",
    [
        "🏢 INDUSTRY",
        "   • Eric Zhu (Alibaba) - AutoGen, QwenPaw",
        "   • Fatma Özcan (Google) - 100x efficiency gains",
        "   • Yunyao Li (Adobe) - Enterprise AI at scale",
        "",
        "🎓 ACADEMIA",
        "   • Eugene Wu (Columbia) - Agentic data environments",
        "   • Juliana Freire (NYU) - Semantic data systems",
        "   • Omar Khattab (MIT) - DSPy, qualitative learning",
        "   • Shreya Shankar (Berkeley→CMU) - What users actually do",
        "",
        "💡 Common Thread Across All Speakers:",
        "   Data systems and agents must co-evolve",
        "   Not data OR agents — data AND agents together"
    ]
)

# ============================================================================
# SLIDE 8: Panel
# ============================================================================
print("Creating Slide 8: Panel...")
slide8 = add_content_slide(
    "Panel Discussion: The Hard Questions",
    [
        "🎤 5:15-6:00 PM Closing Panel",
        "",
        "Panelists: Eric Zhu, Fatma Özcan, Eugene Wu,",
        "           Juliana Freire, Yunyao Li",
        "",
        "Topics:",
        "❓ What abstractions are fundamentally broken?",
        "❓ How do we evaluate compound systems?",
        "❓ What's the right level of human oversight?",
        "❓ Where is production 5 years ahead of research?",
        "❓ What should PhD students work on?",
        "",
        "💬 PLUS: 20 minutes for YOUR questions!",
        "   Hard questions welcome. No polite consensus."
    ]
)

# ============================================================================
# SLIDE 9: What Makes Us Different
# ============================================================================
print("Creating Slide 9: What Makes Us Different...")
slide9 = add_content_slide(
    "What Makes DASHSys Different",
    [
        "1️⃣  Integration, Not Separation",
        "    Papers bridge data systems ↔ AI systems",
        "    Not 'which side wins' but 'how they work together'",
        "",
        "2️⃣  Human-in-the-Loop is Non-Negotiable",
        "    58% of papers include HITL",
        "    Not automation vs. oversight",
        "    But automation WITH oversight",
        "",
        "3️⃣  Production Reality Meets Research Vision",
        "    Industry speakers + academic depth",
        "    Systems track tests ideas on real data",
        "    Gap analysis → research agenda",
        "",
        "🎯 We're not just publishing papers.",
        "   We're defining how data systems evolve for an agentic era."
    ]
)

# ============================================================================
# SLIDE 10: Thank You
# ============================================================================
print("Creating Slide 10: Thank You...")
slide10 = add_title_slide(
    "Thank You!",
    "Let's build the future of data-AI systems — together.\n\n#DASHSys2026"
)

# Add next steps at bottom
left = Inches(2)
top = Inches(5.5)
width = Inches(9.333)
height = Inches(1.5)
text_box = slide10.shapes.add_textbox(left, top, width, height)
tf = text_box.text_frame
tf.text = "📸 Share: #DASHSys2026\n📋 Posters: Coffee breaks\n💬 Questions: Submit anytime\n\nNext: Eric Zhu - 'Surfing the Jagged Frontier'"
tf.paragraphs[0].font.size = Pt(24)
tf.paragraphs[0].font.color.rgb = COLORS['white']
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Save presentation
output_file = "DASHSys-2026-Welcome-Address.pptx"
prs.save(output_file)
print(f"\n✅ PowerPoint created: {output_file}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Ready to present!")
